"""add each type of slide to the presentation
Copyright for portions of PyPresent are held by Steve Canny, 2012, 2013"""

# standard library
from typing import TypeVar, List
import collections.abc             # monkey patch compat/__init__.py

# 3rd party libraries
from pptx.util import Inches, Pt

PT_SIZE = 50
presentation = TypeVar('Presentation')

def add_new_slide(prs: presentation, data: List[tuple], slide_type: str) -> None:
    """Select the type of slide to add and run the function to do so
    :param prs: presentation object
    :param data: list of tuples in (type, contents) format
    :param slide_type: str of the slide type (title, text, image)
    :return: None
    """
    to_function = {'title': add_title, 'text': add_text, 'image': add_image}
    add_type_of_slide = to_function.get(slide_type)
    if not add_type_of_slide:
        print('Error - type of slide does not exist')
    else:
        add_type_of_slide(prs, data)

def add_title(prs: presentation, data: List[tuple]) -> None:
    """Add a title slide
    :param prs: presentation object
    :param data: list of tuples in (type, contents) format
    :return: None
    """
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # title and subtitle
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = '\t'
    subtitle.text = '\t'              # leave blank to fill later

    for data_type, contents in data:
        if data_type == 'title':
            title.text = contents
            title.text_frame.paragraphs[0].font.size = Pt(PT_SIZE)
        elif data_type == 'subtitle':
            subtitle.text = contents

def add_text(prs: presentation, data: List[tuple]) -> None:
    """Add a slide with bullet points
    :param prs: presentation object
    :param data: list of tuples in (type, contents) format
    :return: None
    """
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # title and content

    title = slide.shapes.title
    title.text = '\t'

    body = slide.shapes.placeholders[1]
    body_area = body.text_frame
    default_paragraph_unfilled = True

    for data_type, contents in data:
        if data_type == 'title':
            title.text = contents
            title.text_frame.paragraphs[0].font.size = Pt(PT_SIZE)
        elif data_type == 'bullet':
            if default_paragraph_unfilled:
                body_area.text = contents
                body_area.paragraphs[0].font.size = Pt(PT_SIZE)
                default_paragraph_unfilled = False
                continue
            bullet_point = body_area.add_paragraph()
            bullet_point.text = contents
            bullet_point.level = 0
            bullet_point.font.size = Pt(PT_SIZE)

def add_image(prs: presentation, data: List[tuple]) -> None:
    """Add a slide with image
    :param prs: presentation object
    :param data: list of tuples in (type, path) format
    :return: None
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    in_0_3 = Inches(0.3)
    in_0_5 = Inches(0.5)
    in_1 = Inches(1)
    in_1_25 = Inches(1.25)
    in_1_75 = Inches(1.75)
    in_1_3 = Inches(1.3)
    in_5_5 = Inches(5.5)
    in_9 = Inches(9)
    for data_type, contents in data:
        if data_type == 'title':
            title = slide.shapes.add_textbox(left=in_0_5, top=in_0_3, width=in_9, height=in_1_25)
            title_frame = title.text_frame
            title_frame.text = contents
            title_frame.paragraphs[0].font.size = Pt(PT_SIZE)
        elif data_type == 'picture':
            pic = slide.shapes.add_picture(contents, left=in_0_5, top=in_1_75, height=in_5_5)
