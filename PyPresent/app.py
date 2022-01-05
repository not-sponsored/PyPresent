"""parse arguments and run the main function to create the presentation
Copyright for portions of PyPresent are held by Steve Canny, 2012, 2013"""

# standard library
import sys
import argparse
from typing import TypeVar
import collections.abc             # monkey patch compat/__init__.py

# 3rd party libraries
from pptx import Presentation

# internal imports
import PyPresent.parse_input as parse_input
import PyPresent.add_slide as add_slide
import PyPresent.__init__

Namespace = TypeVar("Namespace")  # type of parser return

def main(passed_arguments: list=[]) -> int:
    """Generate the presentation based on the input options
    :param passed_arguments: list of args including files and options
    :return: int of exit status
    """
    if not passed_arguments:
        passed_arguments = sys.argv[1:]
    args = parser(passed_arguments)

    contents = parse_input.parse_lines(parse_input.read_file(args.in_file))

    prs = Presentation()
    for data in contents:
        slide_type = data.pop()
        add_slide.add_new_slide(prs, data, slide_type)

    prs.save(args.out_file)
    return 0

def parser(args: list) -> Namespace:
    """Return a namespace with the parsed arguments
    :param args: list of arguments
    :return: Namespace of the parsed arguments
    """
    # get the documentation from the package document string
    pypresent_doc_str = PyPresent.__init__.__doc__.split('\n')[2:]
    examples = '\n'.join(pypresent_doc_str)
    parser = argparse.ArgumentParser(
        description='PyPresent - create presentations from existing notes',
        epilog=examples,
        formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument('in_file', help='file to read from')
    parser.add_argument('out_file', help='file to write to')
    return parser.parse_args(args)

if __name__ == '__main__':
    main()
